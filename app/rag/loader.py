"""
Módulo de carga de documentos.
Soporta múltiples formatos: PDF, DOCX, TXT, MD, CSV, XLSX, JSON, HTML, PPTX.
Recorre las carpetas de departamentos y carga todos los archivos disponibles.
"""

import json
from pathlib import Path
from typing import Any, List

from bs4 import BeautifulSoup

try:
    from langchain_core.documents import Document
except ImportError:
    try:
        from langchain.schema import Document  # type: ignore[import-not-found]
    except ImportError:  # pragma: no cover - fallback para entornos sin langchain clásico
        Document = Any  # type: ignore[misc,assignment]

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore[import-not-found]
    except ImportError:  # pragma: no cover - fallback defensivo
        RecursiveCharacterTextSplitter = None  # type: ignore[assignment]

from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    CSVLoader,
    UnstructuredExcelLoader,
)

from app.core.config import (
    DOCUMENTS_DIR,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    SUPPORTED_EXTENSIONS,
    DEPARTMENTS,
)
from app.utils import get_logger

logger = get_logger(__name__)

LOADER_MAP = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
    ".txt": TextLoader,
    ".md": TextLoader,
    ".csv": CSVLoader,
    ".xlsx": UnstructuredExcelLoader,
}


def load_json_document(file_path: str) -> List[Any]:
    """Carga y convierte un archivo JSON a Document de LangChain."""
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    formatted_text = json.dumps(data, indent=2, ensure_ascii=False)
    return [Document(page_content=formatted_text, metadata={"source": file_path})]


def load_html_document(file_path: str) -> List[Any]:
    """Carga y limpia el contenido de un archivo HTML."""
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    for element in soup(["script", "style", "head", "title", "meta"]):
        element.extract()

    text = soup.get_text(separator="\n")
    lines = (line.strip() for line in text.splitlines())
    clean_text = "\n".join(chunk for chunk in lines if chunk)

    return [Document(page_content=clean_text, metadata={"source": file_path})]


def load_pptx_document(file_path: str) -> List[Any]:
    """Carga texto desde una presentación PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_runs = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_runs.append(f"--- Diapositiva {slide_idx} ---\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(text_runs)
        return [Document(page_content=full_text, metadata={"source": file_path})]
    except Exception as e:
        logger.error(f"Error cargando PPTX {file_path}: {e}")
        return []


def load_single_document(file_path: str) -> List[Any]:
    """Carga un único documento según su extensión."""
    ext = Path(file_path).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        logger.warning(f"Extensión no soportada: {ext} en {file_path}")
        return []

    try:
        if ext == ".json":
            docs = load_json_document(file_path)
        elif ext == ".html":
            docs = load_html_document(file_path)
        elif ext == ".pptx":
            docs = load_pptx_document(file_path)
        elif ext in LOADER_MAP:
            loader_class = LOADER_MAP[ext]
            if ext in (".txt", ".md"):
                loader = loader_class(file_path, encoding="utf-8")
            else:
                loader = loader_class(file_path)
            docs = loader.load()
        else:
            logger.warning(f"No hay handler configurado para: {ext}")
            return []

        logger.info(f"✅ Cargado: {file_path} ({len(docs)} documento/s)")
        return docs
    except Exception as e:
        logger.error(f"❌ Error cargando {file_path}: {e}")
        return []


def load_all_documents() -> List[Any]:
    """Recorre todos los departamentos y carga todos los documentos soportados."""
    all_documents = []

    for department in DEPARTMENTS:
        dept_path = DOCUMENTS_DIR / department

        if not dept_path.exists():
            logger.info(f"📁 Creando directorio: {dept_path}")
            dept_path.mkdir(parents=True, exist_ok=True)
            continue

        for file_path in dept_path.rglob("*"):
            if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                docs = load_single_document(str(file_path))

                for doc in docs:
                    doc.metadata["department"] = department
                    doc.metadata["source_file"] = file_path.name

                all_documents.extend(docs)

    logger.info(f"📄 Total de documentos cargados: {len(all_documents)}")
    return all_documents


def split_documents(documents: List[Any]) -> List[Any]:
    """Divide los documentos en chunks para el procesamiento RAG."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = text_splitter.split_documents(documents)
    logger.info(f"✂️ Documentos divididos en {len(chunks)} chunks")
    return chunks


def load_and_split() -> List[Any]:
    """Pipeline completo: carga y divide documentos."""
    documents = load_all_documents()

    if not documents:
        logger.warning("⚠️ No se encontraron documentos para procesar.")
        return []

    return split_documents(documents)
